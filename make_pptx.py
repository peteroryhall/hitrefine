#!/usr/bin/env python3
"""
Compile a folder of 2D interaction PNGs into a PowerPoint presentation.
One slide per ligand with the molecule name as title.
"""

import argparse
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def parse_args():
    p = argparse.ArgumentParser(
        description="Compile ProLIF interaction PNGs into a PowerPoint."
    )
    p.add_argument('-i', '--input',  required=True, help='Directory of PNG files')
    p.add_argument('-o', '--output', required=True, help='Output .pptx file path')
    return p.parse_args()


def main():
    args = parse_args()
    pngs = sorted(Path(args.input).glob('*.png'))
    if not pngs:
        print('No PNG files found.')
        return

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for png in pngs:
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.5))
        tf = txBox.text_frame
        tf.text = png.stem
        tf.paragraphs[0].runs[0].font.size = Pt(10)
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        img = Image.open(str(png))
        img_w, img_h = img.size
        aspect = img_w / img_h
        max_w, max_h = Inches(12.0), Inches(6.7)
        if max_w / aspect <= max_h:
            w, h = max_w, max_w / aspect
        else:
            w, h = max_h * aspect, max_h
        left = (prs.slide_width - w) / 2
        slide.shapes.add_picture(str(png), left, Inches(0.6), w, h)

    prs.save(args.output)
    print(f'Saved {len(pngs)} slides to {args.output}')


if __name__ == '__main__':
    main()
